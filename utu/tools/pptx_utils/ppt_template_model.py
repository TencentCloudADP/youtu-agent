"""
This module is used to define the PPT template pydantic models.
"""

import logging
import traceback
import uuid
import random
import os
from typing import Any, Literal, Tuple, Dict, List, Union, Optional

import requests
import numpy as np
from PIL import Image
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pydantic import BaseModel
from .utils import delete_shape, find_shape_with_name_except, replace_picture_keep_format
from .pyecharts_utils import (
    inject_chinese_font_into_html,
    check_playwright_available,
    playwright_snapshot,
)


class BaseContent(BaseModel):
    content_type: Literal["text", "image", "table"]


class Paragraph(BaseModel):
    text: str
    bullet: bool = False
    level: int = 0


class Item(BaseModel):
    title: str  # be very concise within 3 words, or 4 characters
    content: str  # be very concise within 10 words


class TextContent(BaseContent):
    content_type: Literal["text"] = "text"
    paragraph: Union[List[Paragraph], str]


class BasicImage(BaseModel):
    image_url: str  # absolute url


class ImageContent(BaseContent):
    content_type: Literal["image"] = "image"
    image_url: str  # absolute url
    caption: Optional[str] = None  # be very concise within 20 words


class TableContent(BaseContent):
    content_type: Literal["table"] = "table"
    header: List[str]
    rows: List[List[str]]
    caption: Optional[str] = None  # be very concise within 20 words
    n_rows: int  # no more than 7
    n_cols: int  # no more than 10


class PyEchartsContent(BaseContent):
    content_type: Literal["pyecharts"] = "pyecharts"
    chart_type: str  # e.g., "bar", "line", "pie", "scatter", "radar", etc.
    chart_data: Dict[str, Any]  # chart configuration and data
    width: Optional[str] = "800px"  # chart width
    height: Optional[str] = "600px"  # chart height
    caption: Optional[str] = None  # be very concise within 20 words


# font awesome 4 only
class FontAwesomeIcon(BaseModel):
    icon_name: str


class PageConfig:
    """Configuration loader for page templates from YAML"""

    def __init__(self, config: Dict[str, Any]):
        self.type_map = {}
        self.pages = {}
        self._load_config(config)

    def _load_config(self, config: Dict[str, Any]):
        """Load configuration from YAML config"""

        # Load type_map
        if "type_map" in config:
            for item in config["type_map"]:
                if isinstance(item, dict):
                    for key, value in item.items():
                        # Support both single index and list of indices
                        if isinstance(value, list):
                            self.type_map[key] = value
                        else:
                            self.type_map[key] = [value]
        else:
            raise ValueError("type_map not found in YAML config")

        # Load page configurations (allow both '<type>_page' and '<type>')
        for key, value in config.items():
            if key == "type_map":
                continue

            page_key = key
            self.pages[page_key] = value

    def render(self, slide, page_json: Dict[str, Any], prs):
        """Render slide based on page configuration and data"""
        page_type = page_json.get("type", "")
        logging.info(f"===Rendering page type: {page_type}===")

        # Get page configuration

        page_config = self.pages.get(page_type, {})

        # Render all fields based on their type from YAML config
        for field_name, field_config in page_config.items():
            if field_name == "type" or field_name == "description":
                continue

            field_value = page_json.get(field_name)
            if field_value is None:
                raise ValueError(f"Field '{field_name}' not found in page JSON")

            field_type = field_config.get("type", "str")
            print(f"field_name: {field_name}, field_type: {field_type}")

            if field_type == "str":
                self._render_text_field(slide, field_name, field_value)
            elif field_type == "int":
                int_str = str(field_value)
                self._render_text_field(slide, field_name, int_str)
            elif field_type == "content":
                self._render_content_field(slide, field_name, field_value)
            elif field_type == "content_list":
                self._render_content_list_field(slide, field_name, field_value)
            elif field_type == "item_list":
                self._render_item_list_field(slide, field_name, field_value)
            elif field_type == "str_list":
                self._render_label_list_field(slide, field_name, field_value)
            elif field_type == "image":
                self._render_basic_image_field(slide, field_name, field_value)
            elif field_type == "fontawesome_icon":
                self._render_font_awesome_icon_field(slide, field_name, field_value, prs)
            else:
                raise ValueError(f"Unknown field type: {field_type}")

    def _render_font_awesome_icon_field(self, slide, field_name: str, icon_value, prs):
        logging.info(f"{field_name}: {icon_value}")
        shape = find_shape_with_name_except(slide.shapes, field_name)
        if shape:
            handle_font_awesome_icon(icon_value, shape, slide, prs)

    def _render_basic_image_field(self, slide, field_name: str, image_value):
        logging.info(f"{field_name}: {image_value}")
        shape = find_shape_with_name_except(slide.shapes, field_name)
        image = self._ensure_basic_image_model(image_value)
        if shape:
            handle_image(image.image_url, shape, slide)

    def _render_text_field(self, slide, field_name: str, text_value: str):
        """Render text field"""
        shape = find_shape_with_name_except(slide.shapes, field_name)
        if shape:
            handle_pure_text(text_value, shape, slide)
        else:
            raise ValueError(f"Shape not found for {field_name}")

    def _render_content_field(self, slide, field_name: str, content_value):
        """Render content field"""
        logging.info(f"{field_name}: {content_value}")
        # Use the field name directly to find the shape
        shape = find_shape_with_name_except(slide.shapes, field_name)
        if shape:
            handle_content(self._ensure_content_model(content_value), shape, slide)

    def _render_content_list_field(self, slide, field_name: str, content_values: list):
        """Render list of content fields into <field_name>1, <field_name>2, ..."""
        for i, content_value in enumerate(content_values):
            target_name = f"{field_name}{i + 1}"
            logging.info(f"{target_name}: {content_value}")
            shape = find_shape_with_name_except(slide.shapes, target_name)
            if shape:
                handle_content(self._ensure_content_model(content_value), shape, slide)
            else:
                raise ValueError(f"Shape not found for {target_name}")

    def _render_item_list_field(self, slide, field_name: str, items: list):
        """Render item list field"""
        for ind, item in enumerate(items):
            logging.info(f"{field_name} {ind}: {item}")
            handle_item(self._ensure_item_model(item), ind, slide)

    def _render_label_list_field(self, slide, field_name: str, labels: list):
        """Render label list field"""
        for i, label_text in enumerate(labels):
            label_shape = find_shape_with_name_except(slide.shapes, f"{field_name}{i + 1}")
            if label_shape:
                logging.info(f"{field_name}{i + 1}: {label_text}")
                handle_pure_text(label_text, label_shape, slide)

    def _ensure_basic_image_model(self, image_value: Any) -> BasicImage:
        if isinstance(image_value, BasicImage):
            return image_value

        if not isinstance(image_value, dict):
            raise TypeError("Image field must be dict or BasicImage instance")

        return BasicImage(**image_value)

    def _ensure_content_model(self, content_value: Any) -> BaseContent:
        """Convert dict/list payloads into BaseContent instances."""
        if isinstance(content_value, BaseContent):
            return content_value

        if not isinstance(content_value, dict):
            raise TypeError("Content field must be dict or BaseContent instance")

        content_type = content_value.get("content_type")
        if content_type == "text":
            return TextContent(**content_value)
        if content_type == "image":
            return ImageContent(**content_value)
        if content_type == "table":
            return TableContent(**content_value)
        if content_type == "pyecharts":
            return PyEchartsContent(**content_value)

        raise ValueError(f"Unsupported content type: {content_type}")

    def _ensure_item_model(self, item_value: Any) -> Item:
        """Convert dict payloads into Item instances."""
        if isinstance(item_value, Item):
            return item_value

        if not isinstance(item_value, dict):
            raise TypeError("Item field must be dict or Item instance")

        return Item(**item_value)


def download_image(url, base_dir="."):
    if url.startswith("http"):
        # download the image
        headers = {"Accept": "image/*, */*"}
        response = requests.get(url, headers=headers)
        extension_name = url.split(".")[-1]
        if extension_name not in ["png", "jpg", "jpeg", "gif", "bmp", "webp"]:
            extension_name = "png"
        if response.status_code == 200:
            file_name = f"{base_dir}/{uuid.uuid4()}.{extension_name}"
            with open(file_name, "wb") as f:
                f.write(response.content)
            # get width and height
            image = Image.open(file_name)
            width, height = image.size
            return file_name, width, height
        else:
            raise Exception(f"Failed to download image: {url} {response.status_code}")
    raise Exception(f"Failed to download image: {url}")


# def handle_pure_text(text: str, target_shape, slide):
#     try:
#         text_frame = target_shape.text_frame
#         for paragraph in text_frame.paragraphs:
#             for run in paragraph.runs:
#                 run.text = text
#                 text = ""
#     except Exception as e:
#         logging.error(f"Failed to set text: {text} {e}")
#         traceback.print_exc()


def handle_pure_text(text: str, target_shape, slide):
    try:
        logging.info(f"handle_pure_text: {text} =fill=> {target_shape.text_frame.text}")
        text_frame = target_shape.text_frame
        has_runs = any(paragraph.runs for paragraph in text_frame.paragraphs)

        if not has_runs:
            text_frame.clear()
            paragraph = text_frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = text
        else:
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = text
                    text = ""
    except Exception as e:
        traceback.print_exc()
        logging.error(f"Failed to set text: {text} {e}")


def handle_image(image_url: str, target_shape, slide):
    left, top, width, height = target_shape.left, target_shape.top, target_shape.width, target_shape.height
    
    # Distinguish between remote URL and local relative path
    if image_url and image_url.startswith("http"):
        # Case 1: Remote URL - download the image
        try:
            image_url = image_url
            image_url, image_width, image_height = download_image(image_url)
        except Exception as e:
            logging.warning(f"Failed to download image: {image_url} {e}")
            traceback.print_exc()
            return
    elif image_url and os.path.exists(image_url):
        # Case 2: Local relative path - use image_url directly as local path
        try:
            image_url = image_url
            if not image_url:
                logging.warning(f"Image URL is empty for local image: {image_url}")
                return
            
            # Get image dimensions from local file
            img = Image.open(image_url)
            image_width, image_height = img.size
            logging.info(f"Using local image: {image_url} ({image_width}x{image_height})")
        except Exception as e:
            logging.warning(f"Failed to load local image: {image_url} {e}")
            traceback.print_exc()
            return
    else:
        logging.warning(f"Invalid image configuration: missing required fields")
        return

    # scale the image to fit the placeholder
    scale = min(width / image_width, height / image_height)
    image_width *= scale
    image_height *= scale
    # center the image
    left += (width - image_width) / 2
    top += (height - image_height) / 2

    # get parent
    try:
        parent = target_shape._parent._parent
        # check if parent is shape group
        if parent.shape_type == MSO_SHAPE_TYPE.GROUP:
            # get parent left and top
            parent_left, parent_top = parent.left, parent.top
            slide.shapes.add_picture(image_url, left + parent_left, top + parent_top, image_width, image_height)
        else:
            slide.shapes.add_picture(image_url, left, top, image_width, image_height)
    except AttributeError:
        slide.shapes.add_picture(image_url, left, top, image_width, image_height)

    # remove the placeholder
    delete_shape(target_shape)


def handle_replace_image(image_url: str, target_shape, slide):
    # download image
    try:
        image_url, image_width, image_height = download_image(image_url)
    except Exception as e:
        logging.warning(f"Failed to download image: {image_url} {e}")
        traceback.print_exc()
        return
    replace_picture_keep_format(slide, target_shape, image_url)


def handle_table(table_content: TableContent, target_shape, slide):
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    left, top, width, height = target_shape.left, target_shape.top, target_shape.width, target_shape.height
    table = slide.shapes.add_table(table_content.n_rows + 1, table_content.n_cols, left, top, width, height).table
    for i, header_text in enumerate(table_content.header):
        cell = table.cell(0, i)
        cell.text = header_text
        # Center text horizontally and vertically
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, row in enumerate(table_content.rows):
        for j, cell_text in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = cell_text
            # Center text horizontally and vertically
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    table.auto_fit = True

    # remove the placeholder
    delete_shape(target_shape)


def handle_text_content(text_content: TextContent, target_shape, slide):
    if isinstance(text_content.paragraph, list):
        original_font = None
        for paragraph in target_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                original_font = run.font
                break

        text_frame = target_shape.text_frame

        text_frame.clear()

        for paragraph in text_content.paragraph:
            para = text_frame.add_paragraph()
            para.bullet = paragraph.bullet
            para.level = paragraph.level
            run = para.add_run()
            run.text = paragraph.text
            if original_font:
                run.font.name = original_font.name
                run.font.size = original_font.size
                run.font.bold = original_font.bold
                run.font.italic = original_font.italic
                if original_font.color.type == 1:
                    run.font.color.rgb = original_font.color.rgb
    else:
        handle_pure_text(text_content.paragraph, target_shape, slide)


def handle_content(content: BaseContent, target_shape, slide):
    if content.content_type == "text":
        handle_text_content(content, target_shape, slide)
    elif content.content_type == "image":
        handle_image(content.image_url, target_shape, slide)
    elif content.content_type == "table":
        handle_table(content, target_shape, slide)
    elif content.content_type == "pyecharts":
        handle_pyecharts(content, target_shape, slide)
    else:
        raise ValueError(f"Unsupported content type: {content.content_type}")


def handle_item(item: Item, item_index: int, slide, index_start_from_one=True):
    if index_start_from_one:
        item_index += 1
    item_title_name = f"item_title{item_index}"
    item_content_name = f"item_content{item_index}"

    item_title_shape = find_shape_with_name_except(slide.shapes, item_title_name)
    item_content_shape = find_shape_with_name_except(slide.shapes, item_content_name)

    handle_pure_text(item.title, item_title_shape, slide)
    handle_pure_text(item.content, item_content_shape, slide)


def handle_pyecharts(pyecharts_content: "PyEchartsContent", target_shape, slide):
    """
    Handle PyEcharts chart content by rendering it to an image and inserting into the slide.
    
    Uses Playwright for rendering. Install with: pip install playwright && playwright install chromium
    
    Args:
        pyecharts_content: PyEchartsContent instance containing chart configuration
        target_shape: Target shape placeholder in the slide
        slide: The slide object to add the chart to
    """
    try:
        from pyecharts.charts import (
            Bar, Line, Pie, Scatter, Radar, Funnel, Gauge, 
            WordCloud, HeatMap, Kline, Map, Geo, Graph, Tree,
            Sunburst, Sankey, ThemeRiver, Calendar, Boxplot,
            EffectScatter, Parallel, Polar, Liquid
        )
        from pyecharts import options as opts
    except ImportError as e:
        logging.error(f"Failed to import pyecharts: {e}")
        raise ImportError(
            "PyEcharts support requires 'pyecharts'. "
            "Install with: pip install pyecharts"
        )

    # Map chart type string to chart class
    chart_type_map = {
        "bar": Bar,
        "line": Line,
        "pie": Pie,
        "scatter": Scatter,
        "radar": Radar,
        "funnel": Funnel,
        "gauge": Gauge,
        "wordcloud": WordCloud,
        "heatmap": HeatMap,
        "kline": Kline,
        "map": Map,
        "geo": Geo,
        "graph": Graph,
        "tree": Tree,
        "sunburst": Sunburst,
        "sankey": Sankey,
        "themeriver": ThemeRiver,
        "calendar": Calendar,
        "boxplot": Boxplot,
        "effectscatter": EffectScatter,
        "parallel": Parallel,
        "polar": Polar,
        "liquid": Liquid,
    }

    chart_type = pyecharts_content.chart_type.lower()
    if chart_type not in chart_type_map:
        raise ValueError(
            f"Unsupported chart type: {chart_type}. "
            f"Supported types: {list(chart_type_map.keys())}"
        )

    chart_data = pyecharts_content.chart_data
    width = pyecharts_content.width
    height = pyecharts_content.height

    # Generate unique filenames for temporary files
    temp_html_path = f"{uuid.uuid4()}_pyecharts.html"
    temp_image_path = f"{uuid.uuid4()}_pyecharts.png"

    try:
        # Create chart instance with specified dimensions and transparent background
        # Use a Chinese-compatible font to avoid garbled text on Linux servers
        chart_class = chart_type_map[chart_type]
        chart = chart_class(init_opts=opts.InitOpts(
            width=width, 
            height=height,
            bg_color="rgba(0,0,0,0)",  # Transparent background for better PPT integration
        ))
        
        # Set global font to support Chinese characters
        # Use common fonts that are available on most systems
        chinese_font = "Microsoft YaHei, SimHei, PingFang SC, Noto Sans SC, WenQuanYi Micro Hei, sans-serif"
        chart.set_global_opts(
            title_opts=opts.TitleOpts(
                title_textstyle_opts=opts.TextStyleOpts(font_family=chinese_font)
            ),
            legend_opts=opts.LegendOpts(
                textstyle_opts=opts.TextStyleOpts(font_family=chinese_font),
                item_gap=20,  # Increase gap between legend items to avoid overlap
                orient="vertical",  # Vertical layout to avoid text overlap
                pos_left="right",  # Position legend on the right side of chart
                pos_top="middle",  # Vertically center the legend
            ),
            xaxis_opts=opts.AxisOpts(
                axislabel_opts=opts.LabelOpts(font_family=chinese_font)
            ),
            yaxis_opts=opts.AxisOpts(
                axislabel_opts=opts.LabelOpts(font_family=chinese_font)
            ),
        )

        # Apply chart configuration from chart_data
        if "x_axis" in chart_data:
            chart.add_xaxis(chart_data["x_axis"])
        
        if "y_axis" in chart_data:
            # For bar/line charts, y_axis contains series data
            if isinstance(chart_data["y_axis"], list):
                if chart_type in ["bar", "line"]:
                    for series in chart_data["y_axis"]:
                        if isinstance(series, dict):
                            series_name = series.get("name", "")
                            series_data = series.get("data", [])
                            chart.add_yaxis(
                                series_name, 
                                series_data,
                                label_opts=opts.LabelOpts(font_family=chinese_font)
                            )
                        else:
                            chart.add_yaxis(
                                "", 
                                series,
                                label_opts=opts.LabelOpts(font_family=chinese_font)
                            )
                else:
                    chart.add_yaxis(chart_data["y_axis"])
        
        # For pie chart
        if "data_pair" in chart_data:
            chart.add(
                "", 
                chart_data["data_pair"],
                label_opts=opts.LabelOpts(font_family=chinese_font)
            )
        
        # For radar chart
        if "schema" in chart_data:
            chart.add_schema(schema=chart_data["schema"])
        if "series_data" in chart_data:
            for series in chart_data["series_data"]:
                series_name = series.get("name", "")
                series_values = series.get("value", [])
                chart.add(
                    series_name, 
                    series_values,
                    label_opts=opts.LabelOpts(font_family=chinese_font)
                )
        
        # Apply global options if provided (merge with existing Chinese font settings)
        global_opts_kwargs = {}
        
        if "title" in chart_data:
            global_opts_kwargs['title_opts'] = opts.TitleOpts(
                title=chart_data["title"],
                title_textstyle_opts=opts.TextStyleOpts(font_family=chinese_font)
            )
        
        if "legend" in chart_data:
            legend_config = chart_data["legend"].copy() if isinstance(chart_data["legend"], dict) else {}
            legend_config['textstyle_opts'] = opts.TextStyleOpts(font_family=chinese_font)
            # Set default values to avoid legend overlap if not specified
            legend_config.setdefault('item_gap', 20)  # Gap between legend items
            legend_config.setdefault('orient', 'vertical')  # Vertical layout to avoid text overlap
            legend_config.setdefault('pos_left', 'right')  # Position legend on the right side
            legend_config.setdefault('pos_top', 'middle')  # Vertically center the legend
            global_opts_kwargs['legend_opts'] = opts.LegendOpts(**legend_config)
        
        if global_opts_kwargs:
            chart.set_global_opts(**global_opts_kwargs)

        # Render chart to HTML file first
        chart.render(temp_html_path)
        logging.info(f"PyEcharts chart rendered to HTML: {temp_html_path}")
        
        # Inject web font CSS into the HTML to ensure Chinese characters display correctly
        inject_chinese_font_into_html(temp_html_path, chinese_font)
        
        # Check if Playwright is available
        is_available, error_msg = check_playwright_available()
        if not is_available:
            raise ImportError(error_msg)
        
        # Use Playwright to take snapshot
        try:
            logging.info("Taking snapshot with Playwright...")
            playwright_snapshot(temp_html_path, temp_image_path)
            logging.info(f"PyEcharts chart snapshot saved to: {temp_image_path}")
        except Exception as e:
            error_str = str(e).lower()
            # Provide helpful error message for common issues
            if "playwright install" in error_str or "executable doesn't exist" in error_str:
                raise RuntimeError(
                    f"Playwright browser not installed. Please run: playwright install chromium\n"
                    f"Original error: {e}"
                )
            elif "no usable sandbox" in error_str or "running as root" in error_str:
                raise RuntimeError(
                    f"Playwright sandbox error. Try running with --no-sandbox flag or as non-root user.\n"
                    f"Original error: {e}"
                )
            elif "cannot open display" in error_str or "display" in error_str:
                raise RuntimeError(
                    f"Display not available. Ensure headless mode is enabled or X server is running.\n"
                    f"Original error: {e}"
                )
            else:
                raise RuntimeError(
                    f"Failed to take snapshot with Playwright: {e}\n"
                    "Troubleshooting tips:\n"
                    "  1. Ensure browser is installed: playwright install chromium\n"
                    "  2. Check system dependencies: playwright install-deps chromium\n"
                    "  3. On Linux, you may need additional libraries - see Playwright docs"
                )
        
        # Use existing handle_image function to insert the chart image
        handle_image(temp_image_path, target_shape, slide)
        
    except Exception as e:
        logging.error(f"Failed to render PyEcharts chart: {e}")
        traceback.print_exc()
        raise
    finally:
        # Clean up temporary files
        for temp_file in [temp_html_path, temp_image_path]:
            try:
                if os.path.exists(temp_file):
                    os.remove(temp_file)
            except Exception as cleanup_error:
                logging.warning(f"Failed to clean up temporary file {temp_file}: {cleanup_error}")


def handle_font_awesome_icon(icon: FontAwesomeIcon, target_shape, slide, prs):
    from .fa_to_codepoint import fa_2_codepoint
    from .utils import get_fill_rgb
    
    # random icons: [star, "!", smile face, "*", thumb up, "v"]
    RANDOM_ICONS = [61445, 61546, 61720, 61545, 61796, 61754]

    codepoint = fa_2_codepoint.get(icon["icon_name"])
    if codepoint is None:
        codepoint = random.choice(RANDOM_ICONS)
    
    target_color = get_fill_rgb(target_shape, prs)
    if not target_color:
        logging.warning("cannot get target color of icon, use black by default")
        target_color = (0, 0, 0)
    
    img = Image.open(f'fa_icons/{codepoint}.png')
    
    # replace black in the image with the target_color
    img_array = np.array(img)
    # 分离RGB和Alpha通道
    rgb = img_array[:, :, :3]
    alpha = img_array[:, :, 3] if img_array.shape[2] == 4 else None
    
    # 创建黑色像素掩码
    black_mask = np.all(rgb <= 10, axis=2)
    
    # 替换RGB通道中的黑色
    rgb[black_mask] = target_color
    
    # 重新组合通道
    if alpha is not None:
        result_array = np.dstack((rgb, alpha))
    else:
        result_array = rgb
    result_img = Image.fromarray(result_array)
    
    # save to temp file
    temp_file = f"{uuid.uuid4()}.png"
    result_img.save(temp_file)
    
    handle_image(temp_file, target_shape, slide)
