import copy
import json
import logging
import random
import traceback
import uuid
import requests
import matplotlib
from PIL import Image
from pydantic import BaseModel
from typing import Any, Literal
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation


class BaseContent(BaseModel):
    content_type: Literal["text", "image", "table"]


class Paragraph(BaseModel):
    text: str
    bullet: bool = False
    level: int = 0


class Item(BaseModel):
    title: str  # be very concise within 3 words, or 4 characters
    content: str  # be very concise within 10 words


class TextContent(BaseContent):
    content_type: Literal["text"] = "text"
    paragraph: list[Paragraph] | str


class BasicImage(BaseModel):
    image_url: str  # absolute url


class ImageContent(BaseContent):
    content_type: Literal["image"] = "image"
    image_url: str  # absolute url
    caption: str | None = None  # be very concise within 20 words


class TableContent(BaseContent):
    content_type: Literal["table"] = "table"
    header: list[str]
    rows: list[list[str]]
    caption: str | None = None  # be very concise within 20 words
    n_rows: int  # no more than 7
    n_cols: int  # no more than 10


class PageConfig:
    """Configuration loader for page templates from YAML"""

    def __init__(self, config: dict[str, Any]):
        self.type_map = {}
        self.pages = {}
        self._load_config(config)

    def _load_config(self, config: dict[str, Any]):
        """Load configuration from YAML config"""

        # Load type_map
        if "type_map" in config:
            for item in config["type_map"]:
                if isinstance(item, dict):
                    for key, value in item.items():
                        self.type_map[key] = value
        else:
            raise ValueError("type_map not found in YAML config")

        # Load page configurations (allow both '<type>_page' and '<type>')
        for key, value in config.items():
            if key == "type_map":
                continue

            page_key = key if key.endswith("_page") else f"{key}_page"
            self.pages[page_key] = value

    def render(self, slide, page_json: dict[str, Any]):
        """Render slide based on page configuration and data"""
        page_type = page_json.get("type", "")
        logging.info(f"===Rendering page type: {page_type}===")

        # Get page configuration
        page_config = self.pages.get(f"{page_type}_page", {})

        # Render all fields based on their type from YAML config
        for field_name, field_config in page_config.items():
            if field_name == "type" or field_name == "description":
                continue

            field_value = page_json.get(field_name)
            if field_value is None:
                continue

            field_type = field_config.get("type", "str")

            if field_type == "str":
                self._render_text_field(slide, field_name, field_value)
            elif field_type == "int":
                int_str = str(field_value)
                self._render_text_field(slide, field_name, int_str)
            elif field_type == "content":
                self._render_content_field(slide, field_name, field_value)
            elif field_type == "content_list":
                self._render_content_list_field(slide, field_name, field_value)
            elif field_type == "item_list":
                self._render_item_list_field(slide, field_name, field_value)
            elif field_type == "str_list":
                self._render_label_list_field(slide, field_name, field_value)
            elif field_type == "image":
                self._render_basic_image_field(slide, field_name, field_value)
            else:
                logging.warning(f"Unknown field type: {field_type}")

    def _render_basic_image_field(self, slide, field_name: str, image_value):
        logging.info(f"{field_name}: {image_value}")
        shape = find_shape_with_name_except(slide.shapes, field_name)
        image = self._ensure_basic_image_model(image_value)
        if shape:
            handle_image(image.image_url, shape, slide)

    def _render_text_field(self, slide, field_name: str, text_value: str):
        """Render text field"""
        logging.info(f"{field_name}: {text_value}")
        shape = find_shape_with_name_except(slide.shapes, field_name)
        if shape:
            handle_pure_text(text_value, shape, slide)

    def _render_content_field(self, slide, field_name: str, content_value):
        """Render content field"""
        logging.info(f"{field_name}: {content_value}")
        # Use the field name directly to find the shape
        shape = find_shape_with_name_except(slide.shapes, field_name)
        if shape:
            handle_content(self._ensure_content_model(content_value), shape, slide)

    def _render_content_list_field(self, slide, field_name: str, content_values: list):
        """Render list of content fields into <field_name>1, <field_name>2, ..."""
        for i, content_value in enumerate(content_values):
            target_name = f"{field_name}{i + 1}"
            logging.info(f"{target_name}: {content_value}")
            shape = find_shape_with_name_except(slide.shapes, target_name)
            if shape:
                handle_content(self._ensure_content_model(content_value), shape, slide)
            else:
                logging.warning(f"Shape not found for {target_name}")

    def _render_item_list_field(self, slide, field_name: str, items: list):
        """Render item list field"""
        for ind, item in enumerate(items):
            logging.info(f"{field_name} {ind}: {item}")
            handle_item(self._ensure_item_model(item), ind, slide)

    def _render_label_list_field(self, slide, field_name: str, labels: list):
        """Render label list field"""
        for i, label_text in enumerate(labels):
            label_shape = find_shape_with_name_except(slide.shapes, f"label{i + 1}")
            if label_shape:
                logging.info(f"label{i + 1}: {label_text}")
                handle_pure_text(label_text, label_shape, slide)

    def _ensure_basic_image_model(self, image_value: Any) -> BasicImage:
        if isinstance(image_value, BasicImage):
            return image_value

        if not isinstance(image_value, dict):
            raise TypeError("Image field must be dict or BasicImage instance")

        return BasicImage(**image_value)

    def _ensure_content_model(self, content_value: Any) -> BaseContent:
        """Convert dict/list payloads into BaseContent instances."""
        if isinstance(content_value, BaseContent):
            return content_value

        if not isinstance(content_value, dict):
            raise TypeError("Content field must be dict or BaseContent instance")

        content_type = content_value.get("content_type")
        if content_type == "text":
            return TextContent(**content_value)
        if content_type == "image":
            return ImageContent(**content_value)
        if content_type == "table":
            return TableContent(**content_value)

        raise ValueError(f"Unsupported content type: {content_type}")

    def _ensure_item_model(self, item_value: Any) -> Item:
        """Convert dict payloads into Item instances."""
        if isinstance(item_value, Item):
            return item_value

        if not isinstance(item_value, dict):
            raise TypeError("Item field must be dict or Item instance")

        return Item(**item_value)


def download_image(url, base_dir="."):
    if url.startswith("http"):
        # download the image
        headers = {"Accept": "image/*, */*"}
        response = requests.get(url, headers=headers)
        extension_name = url.split(".")[-1]
        if extension_name not in ["png", "jpg", "jpeg", "gif", "bmp", "webp"]:
            extension_name = "png"
        if response.status_code == 200:
            file_name = f"{base_dir}/{uuid.uuid4()}.{extension_name}"
            with open(file_name, "wb") as f:
                f.write(response.content)
            # get width and height
            image = Image.open(file_name)
            width, height = image.size
            return file_name, width, height
        else:
            raise Exception(f"Failed to download image: {url} {response.status_code}")
    raise Exception(f"Failed to download image: {url}")


# def handle_pure_text(text: str, target_shape, slide):
#     try:
#         text_frame = target_shape.text_frame
#         for paragraph in text_frame.paragraphs:
#             for run in paragraph.runs:
#                 run.text = text
#                 text = ""
#     except Exception as e:
#         logging.error(f"Failed to set text: {text} {e}")
#         traceback.print_exc()


def handle_pure_text(text: str, target_shape, slide):
    try:
        text_frame = target_shape.text_frame
        has_runs = any(paragraph.runs for paragraph in text_frame.paragraphs)

        if not has_runs:
            text_frame.clear()
            paragraph = text_frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = text
            return

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = text
                text = ""
    except Exception as e:
        logging.error(f"Failed to set text: {text} {e}")


def handle_image(image_url: str, target_shape, slide):
    left, top, width, height = target_shape.left, target_shape.top, target_shape.width, target_shape.height
    try:
        image_url, image_width, image_height = download_image(image_url)
    except Exception as e:
        logging.warning(f"Failed to download image: {image_url} {e}")
        traceback.print_exc()
        return

    # scale the image to fit the placeholder
    scale = min(width / image_width, height / image_height)
    image_width *= scale
    image_height *= scale
    # center the image
    left += (width - image_width) / 2
    top += (height - image_height) / 2

    # get parent
    try:
        parent = target_shape._parent._parent
        # check if parent is shape group
        if parent.shape_type == MSO_SHAPE_TYPE.GROUP:
            # get parent left and top
            parent_left, parent_top = parent.left, parent.top
            slide.shapes.add_picture(image_url, left + parent_left, top + parent_top, image_width, image_height)
        else:
            slide.shapes.add_picture(image_url, left, top, image_width, image_height)
    except AttributeError:
        slide.shapes.add_picture(image_url, left, top, image_width, image_height)

    # remove the placeholder
    delete_shape(target_shape)


def handle_replace_image(image_url: str, target_shape, slide):
    # download image
    try:
        image_url, image_width, image_height = download_image(image_url)
    except Exception as e:
        logging.warning(f"Failed to download image: {image_url} {e}")
        traceback.print_exc()
        return
    replace_picture_keep_format(slide, target_shape, image_url)


def handle_table(table_content: TableContent, target_shape, slide):
    left, top, width, height = target_shape.left, target_shape.top, target_shape.width, target_shape.height
    table = slide.shapes.add_table(table_content.n_rows + 1, table_content.n_cols, left, top, width, height).table
    for i, header_text in enumerate(table_content.header):
        table.cell(0, i).text = header_text
    for i, row in enumerate(table_content.rows):
        for j, cell_text in enumerate(row):
            table.cell(i + 1, j).text = cell_text
    table.auto_fit = True

    # remove the placeholder
    delete_shape(target_shape)


def handle_text_content(text_content: TextContent, target_shape, slide):
    if isinstance(text_content.paragraph, list):
        original_font = None
        for paragraph in target_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                original_font = run.font
                break

        text_frame = target_shape.text_frame

        text_frame.clear()

        for paragraph in text_content.paragraph:
            para = text_frame.add_paragraph()
            para.bullet = paragraph.bullet
            para.level = paragraph.level
            run = para.add_run()
            run.text = paragraph.text
            if original_font:
                run.font.name = original_font.name
                run.font.size = original_font.size
                run.font.bold = original_font.bold
                run.font.italic = original_font.italic
                if original_font.color.type == 1:
                    run.font.color.rgb = original_font.color.rgb
    else:
        handle_pure_text(text_content.paragraph, target_shape, slide)


def handle_content(content: BaseContent, target_shape, slide):
    if content.content_type == "text":
        handle_text_content(content, target_shape, slide)
    elif content.content_type == "image":
        handle_image(content.image_url, target_shape, slide)
    elif content.content_type == "table":
        handle_table(content, target_shape, slide)
    else:
        raise ValueError(f"Unsupported content type: {content.content_type}")


def handle_item(item: Item, item_index: int, slide, index_start_from_one=True):
    if index_start_from_one:
        item_index += 1
    item_title_name = f"item_title{item_index}"
    item_content_name = f"item_content{item_index}"

    item_title_shape = find_shape_with_name_except(slide.shapes, item_title_name)
    item_content_shape = find_shape_with_name_except(slide.shapes, item_content_name)

    handle_pure_text(item.title, item_title_shape, slide)
    handle_pure_text(item.content, item_content_shape, slide)


# rgb colors from a color scheme
_color_palette = [
    (178, 34, 34),  # Brick Red
    (46, 139, 87),  # Sea Green
    (70, 130, 180),  # Steel Blue
    (210, 180, 140),  # Tan
    (147, 112, 219),  # Medium Purple
    (255, 165, 0),  # Orange (desaturated)
    (72, 209, 204),  # Medium Turquoise
    (205, 92, 92),  # Indian Red
    (106, 90, 205),  # Slate Blue
    (238, 130, 238),  # Violet
    (60, 179, 113),  # Medium Sea Green
    (100, 149, 237),  # Cornflower Blue
    (218, 165, 32),  # Goldenrod
    (199, 21, 133),  # Medium Violet Red
    (65, 105, 225),  # Royal Blue
]


def inspect_ppt(prs):
    """
    Inspect the given presentation.
    """
    for slide in prs.slides:
        inspect_slide(slide)


def inspect_slide(slide):
    """
    Inspect the given slide layout.
    """

    def _inspect_shape_list(shapes, indent=0):
        for shape in shapes:
            print(" " * indent + shape.name)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _inspect_shape_list(shape.shapes, indent + 2)

    _inspect_shape_list(slide.shapes)


def to_svg(slide, prs, svg_filename="test.svg"):
    """
    Convert the given slide layout to an SVG file.
    """

    def _to_svg_box(shapes, svg_box):
        for shape in shapes:
            svg_box.append(
                {
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                    "shape_type": shape.shape_type,
                }
            )
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                _to_svg_box(shape.shapes, svg_box)

    def render_svg(svg_box):
        width = prs.slide_width.inches
        height = prs.slide_height.inches
        fig, ax = matplotlib.pyplot.subplots(figsize=(width, height))
        ax.set_xlim(0, width)
        ax.set_ylim(height, 0)
        for box in svg_box:
            left, top, width, height = box["left"].inches, box["top"].inches, box["width"].inches, box["height"].inches
            picked_color = random.choice(_color_palette)
            color = (picked_color[0] / 255, picked_color[1] / 255, picked_color[2] / 255, 0.3)
            ax.add_patch(matplotlib.patches.Rectangle((left, top), width, height, color=color))
        ax.set_axis_off()
        matplotlib.pyplot.subplots_adjust(left=0, right=1, top=1, bottom=0)
        matplotlib.pyplot.margins(0)
        ax.set_xmargin(0)
        ax.set_ymargin(0)
        fig.savefig(svg_filename, bbox_inches="tight", pad_inches=0)

    svg_box = []
    _to_svg_box(slide.shapes, svg_box)
    render_svg(svg_box)

    return svg_box


def delete_shape(shape):
    """
    Delete the given shape.
    """
    parent = shape.element.getparent()
    parent.remove(shape.element)


def find_shape_with_name(shapes, name, depth=0):
    """
    Find the shape with the given name in the given shapes.
    """
    if depth == 0:
        logging.info(f"Finding shape with name: {name}")
    for shape in shapes:
        if shape.name == name:
            return shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            found = find_shape_with_name(shape.shapes, name, depth + 1)
            if found:
                return found
    return None


def find_shape_with_name_except(shapes, name, depth=0):
    """
    Find the shape with the given name in the given shapes, except the shape with the given name.
    """
    if depth == 0:
        logging.info(f"Finding shape with name: {name}")
    for shape in shapes:
        if shape.name == name:
            return shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            found = find_shape_with_name(shape.shapes, name, depth + 1)
            if found:
                return found
    raise Exception(f"Shape with name {name} not found")


def duplicate_slide(prs, slide):
    slide_layout = slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    for shape in slide.shapes:
        el = shape.element
        new_el = copy.deepcopy(el)

        # 处理图片 - 使用 python-pptx 内置命名空间
        try:
            blips = new_el.xpath(".//a:blip[@r:embed]")

            for blip in blips:
                old_rId = blip.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")

                if old_rId:
                    # 获取原始图片
                    old_image_part = slide.part.related_part(old_rId)

                    # 在新幻灯片中建立关系
                    new_rId = new_slide.part.relate_to(
                        old_image_part, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
                    )

                    # 更新 rId
                    blip.set("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed", new_rId)

        except (KeyError, AttributeError):
            pass

        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    return new_slide


def delete_slide_range(prs, index_range):
    """delete slides in the given index range
    Args:
        prs: Presentation object
        index_range: range of slide indices (0-based)
    """
    for index in reversed(index_range):
        delete_slide(prs, index)


def delete_slide(prs, index):
    """delete slide at the given index
    Args:
        prs: Presentation object
        index: slide index (0-based)

    Raises:
        IndexError: when index out of range
    """
    if index < 0 or index >= len(prs.slides):
        raise IndexError(f"Slide index {index} out of range (0-{len(prs.slides) - 1})")

    xml_slides = prs.slides._sldIdLst
    xml_slides.remove(xml_slides[index])


def move_slide(prs, old_index, new_index):
    """move slide from old_index to new_index"""
    xml_slides = prs.slides._sldIdLst
    # get the element to move
    slide_element = xml_slides[old_index]
    # remove the element from old position
    xml_slides.remove(slide_element)
    # insert to new position
    xml_slides.insert(new_index, slide_element)
    return prs


def replace_picture_keep_format(slide, shape_index, new_image_path):
    shape = slide.shapes[shape_index]

    if shape.shape_type != 13:
        raise ValueError("Target shape is not a picture")

    img_id = shape._element.blip_rEmbed
    image_part = shape.part.related_part(img_id)

    with open(new_image_path, "rb") as f:
        image_part._blob = f.read()

    return shape


def fill_template_with_yaml_config(template_path, output_path, json_data, yaml_config: dict[str, Any]):
    page_config = PageConfig(yaml_config)
    prs = Presentation(template_path)
    data = json.loads(json_data)
    slides_data = data.get("slides", [])

    if not isinstance(slides_data, list):
        raise ValueError("JSON data must contain a 'slides' list")

    for slide_data in slides_data:
        slide_type = slide_data.get("type")
        if not slide_type:
            logging.warning("Skipped slide without type definition: %s", slide_data)
            continue

        template_index = page_config.type_map.get(slide_type)
        if template_index is None or template_index >= len(prs.slides):
            logging.warning("No template found for slide type '%s'", slide_type)
            continue

        template_slide = prs.slides[template_index]
        if slide_type in ("title", "acknowledgement"):
            target_slide = template_slide
        else:
            target_slide = duplicate_slide(prs, template_slide)

        page_config.render(target_slide, slide_data)

    delete_slide_range(prs, range(2, 12))
    delete_slide(prs, 0)
    move_slide(prs, 1, len(prs.slides) - 1)
    prs.save(output_path)
